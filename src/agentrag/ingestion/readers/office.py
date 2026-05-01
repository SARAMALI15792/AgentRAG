"""Office file readers: .xlsx, .pptx, .csv."""

from __future__ import annotations

import csv
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def read_xlsx(path: Path) -> str:
    """Extract text from an Excel workbook, all sheets, headers preserved."""
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "openpyxl is required for .xlsx support. "
            "Install it with: pip install agentrag[office]"
        ) from None

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                parts.append(line)
    wb.close()
    text = "\n".join(parts)
    if not text.strip():
        raise ValueError(f"No text content extracted from {path}.")
    return text


def read_pptx(path: Path) -> str:
    """Extract text from a PowerPoint presentation, slides and speaker notes."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for .pptx support. "
            "Install it with: pip install agentrag[office]"
        ) from None

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"Slide {i}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"Speaker notes: {notes_text}")
    text = "\n".join(parts)
    if not text.strip():
        raise ValueError(f"No text content extracted from {path}.")
    return text


def read_csv(path: Path) -> str:
    """Extract text from a CSV file: header row followed by data rows."""
    lines: list[str] = []
    with open(path, encoding="utf-8", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            line = ", ".join(row).strip()
            if line:
                lines.append(line)
    text = "\n".join(lines)
    if not text.strip():
        raise ValueError(f"No text content extracted from {path}.")
    return text


# Register all office readers
from agentrag.ingestion.reader_registry import register  # noqa: E402

register([".xlsx"], read_xlsx)
register([".pptx"], read_pptx)
register([".csv"], read_csv)
